"""
Document loader: PDF, Word, Excel, PowerPoint, CSV, TXT, URLs, YouTube, and any video.
YouTube handling now supports:
  - Any language (English, Hindi, Arabic, etc.)
  - Videos without subtitles (using Whisper AI transcription)
  - Auto-detection of spoken language

Any video URL (Vimeo, Dailymotion, etc.) is also supported via yt-dlp + Whisper.
"""
import logging
import re
import os
import tempfile
from pathlib import Path
import requests
from langchain_core.documents import Document

logger = logging.getLogger(__name__)

_MIN_TEXT_LENGTH = 200

# ── Unit normalisation ────────────────────────────────────────────────────────
def _normalize_units(text: str) -> str:
    text = re.sub(r'/\s*(\d+)\s*hr[s]?', r' per \1 hours', text, flags=re.IGNORECASE)
    text = re.sub(r'/\s*hour[s]?\b',     ' per 1 hour',    text, flags=re.IGNORECASE)
    text = re.sub(r'/\s*hr[s]?\b',       ' per 1 hour',    text, flags=re.IGNORECASE)
    text = re.sub(r'/\s*day[s]?\b',      ' per 1 day',     text, flags=re.IGNORECASE)
    text = re.sub(r'/\s*km\b',           ' per 1 km',      text, flags=re.IGNORECASE)
    text = re.sub(r'/\s*night[s]?\b',    ' per 1 night',   text, flags=re.IGNORECASE)
    return text

# ── URL helpers ───────────────────────────────────────────────────────────────
def extract_urls(text: str) -> list[str]:
    return re.findall(r'https?://[^\s\'"<>]+', text)

def is_youtube_url(url: str) -> bool:
    return bool(re.search(r'(youtube\.com/watch|youtu\.be/)', url))

# ------------------------------------------------------------------------------
# ADVANCED YOUTUBE TRANSCRIPT EXTRACTION (with Whisper fallback)
# ------------------------------------------------------------------------------

def _get_youtube_transcript_with_whisper_fallback(url: str) -> tuple[str, dict]:
    """
    Get transcript from YouTube video.
    Returns: (transcript_text, metadata)
    """
    from youtube_transcript_api import YouTubeTranscriptApi
    from youtube_transcript_api._errors import NoTranscriptFound, TranscriptsDisabled
    
    match = re.search(r'(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})', url)
    if not match:
        return "Invalid YouTube URL format.", {"error": "invalid_url"}
    video_id = match.group(1)
    
    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        try:
            transcript = transcript_list.find_manually_created_transcript()
            transcript_text = " ".join(entry["text"] for entry in transcript.fetch())
            language_code = transcript.language_code
            logger.info(f"Using manual transcript in {language_code} for {video_id}")
            return transcript_text, {"source_type": "youtube_manual", "language": language_code}
        except:
            pass
        try:
            transcript = transcript_list.find_generated_transcript()
            transcript_text = " ".join(entry["text"] for entry in transcript.fetch())
            language_code = transcript.language_code
            logger.info(f"Using auto-generated transcript in {language_code} for {video_id}")
            return transcript_text, {"source_type": "youtube_auto", "language": language_code}
        except:
            pass
        for transcript in transcript_list:
            transcript_text = " ".join(entry["text"] for entry in transcript.fetch())
            logger.info(f"Using fallback transcript in {transcript.language_code} for {video_id}")
            return transcript_text, {"source_type": "youtube_fallback", "language": transcript.language_code}
    except (NoTranscriptFound, TranscriptsDisabled) as e:
        logger.info(f"No transcript available for {video_id}, falling back to Whisper transcription.")
        return _transcribe_youtube_audio(url, video_id)
    except Exception as e:
        logger.warning(f"Unexpected error getting transcript: {e}")
        return _transcribe_youtube_audio(url, video_id)


def _transcribe_youtube_audio(url: str, video_id: str) -> tuple[str, dict]:
    """Download audio from YouTube and transcribe using Whisper."""
    import whisper
    import yt_dlp
    
    temp_audio_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            temp_audio_path = tmp.name
        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': temp_audio_path.replace('.mp3', ''),
            'quiet': True,
            'no_warnings': True,
        }
        logger.info(f"Downloading audio from {url}...")
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        actual_path = temp_audio_path.replace('.mp3', '.mp3')
        if not os.path.exists(actual_path):
            import glob
            mp3_files = glob.glob(temp_audio_path.replace('.mp3', '') + "*.mp3")
            if mp3_files:
                actual_path = mp3_files[0]
            else:
                raise Exception("Could not find downloaded audio file")
        whisper_model = _get_whisper_model()
        result = whisper_model.transcribe(actual_path, task="transcribe")
        transcript_text = result["text"].strip()
        detected_language = result["language"]
        if not transcript_text:
            raise Exception("Whisper returned empty transcript")
        logger.info(f"Whisper transcription complete. Detected language: {detected_language}")
        return transcript_text, {"source_type": "whisper_transcription", "language": detected_language}
    except Exception as e:
        logger.error(f"Whisper transcription failed: {e}")
        return f"Could not transcribe audio: {str(e)}", {"error": str(e)}
    finally:
        if temp_audio_path and os.path.exists(temp_audio_path):
            os.unlink(temp_audio_path)
        for ext in ['.mp3', '.temp', '.part']:
            path = temp_audio_path.replace('.mp3', '') + ext if temp_audio_path else None
            if path and os.path.exists(path):
                os.unlink(path)


# Global Whisper model
_whisper_model = None
_whisper_lock = None

def _get_whisper_model():
    global _whisper_model, _whisper_lock
    import asyncio
    if _whisper_lock is None:
        _whisper_lock = asyncio.Lock()
    if _whisper_model is None:
        import whisper
        logger.info("Loading Whisper model (base) - first time may take a moment...")
        _whisper_model = whisper.load_model("base")
        logger.info("Whisper model loaded successfully.")
    return _whisper_model


def _load_youtube(url: str) -> list[Document]:
    """Main entry point for YouTube video processing."""
    transcript_text, metadata = _get_youtube_transcript_with_whisper_fallback(url)
    source_type = metadata.get("source_type", "unknown")
    language = metadata.get("language", "unknown")
    return [Document(
        page_content=f"YouTube Video: {url}\nLanguage: {language}\nTranscript source: {source_type}\n\nTranscript:\n{transcript_text}",
        metadata={
            "source": url,
            "type": "youtube",
            "source_type": source_type,
            "language": language,
            "video_id": re.search(r'(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})', url).group(1) if re.search(r'(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})', url) else None,
        },
    )]


# ------------------------------------------------------------------------------
# GENERIC VIDEO SUPPORT (any video URL)
# ------------------------------------------------------------------------------
def _load_generic_video(url: str) -> list[Document]:
    """
    Download audio from any video URL (Vimeo, Dailymotion, etc.) using yt-dlp,
    transcribe with Whisper, and return Document.
    """
    import whisper
    import yt_dlp
    import tempfile
    import os
    import re

    temp_audio_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            temp_audio_path = tmp.name

        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': temp_audio_path.replace('.mp3', ''),
            'quiet': True,
            'no_warnings': True,
        }
        logger.info(f"Downloading audio from generic video URL: {url}")
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])

        actual_path = temp_audio_path.replace('.mp3', '.mp3')
        if not os.path.exists(actual_path):
            import glob
            mp3_files = glob.glob(temp_audio_path.replace('.mp3', '') + "*.mp3")
            if mp3_files:
                actual_path = mp3_files[0]
            else:
                raise Exception("Could not locate downloaded audio file")

        whisper_model = _get_whisper_model()
        result = whisper_model.transcribe(actual_path, task="transcribe")
        transcript_text = result["text"].strip()
        detected_language = result["language"]

        if not transcript_text:
            raise Exception("Whisper returned empty transcript")

        logger.info(f"Generic video transcribed, language: {detected_language}")
        return [Document(
            page_content=f"Video URL: {url}\nLanguage: {detected_language}\nTranscription source: Whisper\n\nTranscript:\n{transcript_text}",
            metadata={
                "source": url,
                "type": "video",
                "source_type": "generic_video",
                "language": detected_language,
            },
        )]
    except Exception as e:
        logger.error(f"Generic video transcription failed for {url}: {e}")
        return [Document(
            page_content=f"Could not transcribe video from {url}. Error: {e}",
            metadata={"source": url, "type": "video", "error": str(e)},
        )]
    finally:
        if temp_audio_path and os.path.exists(temp_audio_path):
            os.unlink(temp_audio_path)
        for ext in ['.mp3', '.temp', '.part']:
            path = temp_audio_path.replace('.mp3', '') + ext if temp_audio_path else None
            if path and os.path.exists(path):
                os.unlink(path)


# ------------------------------------------------------------------------------
# WEBPAGE URL EXTRACTION (unchanged from previous version)
# ------------------------------------------------------------------------------
def load_url_advanced(url: str) -> list[Document]:
    """Extract clean content from any URL using multiple strategies."""
    try:
        jina_url = f"https://r.jina.ai/{url}"
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        resp = requests.get(jina_url, headers=headers, timeout=15)
        if resp.status_code == 200:
            text = resp.text
            lines = text.split('\n')
            title = lines[0].strip() if lines else url
            content = '\n'.join(lines[1:]) if len(lines) > 1 else text
            return [Document(
                page_content=f"URL: {url}\nTitle: {title}\n\n{content}",
                metadata={"source": url, "extraction": "jina_reader", "title": title}
            )]
    except Exception as e:
        logger.warning("Jina Reader failed for %s: %s", url, e)

    try:
        from readability import Document as ReadabilityDoc
        resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=30)
        resp.raise_for_status()
        readable = ReadabilityDoc(resp.text)
        title = readable.title()
        content = readable.summary()
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if len(text) > 200:
            return [Document(
                page_content=f"URL: {url}\nTitle: {title}\n\n{text}",
                metadata={"source": url, "extraction": "readability"}
            )]
    except Exception as e:
        logger.warning("Readability fallback failed: %s", e)

    return _load_webpage(url)


def _load_webpage(url: str) -> list[Document]:
    text = None
    title = url
    try:
        import trafilatura
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        downloaded = trafilatura.fetch_url(url, headers=headers, timeout=30)
        if downloaded:
            text = trafilatura.extract(downloaded, include_comments=False, include_tables=True)
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(downloaded, "html.parser")
                if soup.title and soup.title.string:
                    title = soup.title.string.strip()
            except Exception:
                pass
    except Exception as e:
        logger.warning("trafilatura failed: %s", e)

    if not text or len(text.strip()) < _MIN_TEXT_LENGTH:
        try:
            from bs4 import BeautifulSoup
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=30)
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                tag.decompose()
            if soup.title and soup.title.string:
                title = soup.title.string.strip()
            bs_text = re.sub(r'\n{3,}', '\n\n', soup.get_text(separator="\n", strip=True))
            if len(bs_text.strip()) > len((text or "").strip()):
                text = bs_text
        except Exception as e:
            logger.warning("requests+bs4 fallback failed: %s", e)

    if not text or len(text.strip()) < 150:
        return [Document(
            page_content=(
                f"⚠️ Could not extract meaningful content from: {url}\n\n"
                "This page likely requires login, JavaScript rendering, or blocks bots."
            ),
            metadata={"source": url, "type": "webpage", "error": "insufficient_content"},
        )]

    return [Document(
        page_content=f"Web Page: {title}\nURL: {url}\n\n{text}",
        metadata={"source": url, "type": "webpage", "title": title},
    )]


def load_url(url: str) -> list[Document]:
    """Main entry point for URL loading – supports webpages and any video URL."""
    if is_youtube_url(url):
        return _load_youtube(url)
    # Check if it looks like a video URL (common patterns)
    video_patterns = [r'(vimeo\.com)', r'(dailymotion\.com)', r'(twitch\.tv)', r'(facebook\.com/watch)', r'(tiktok\.com)']
    if any(re.search(p, url) for p in video_patterns):
        return _load_generic_video(url)
    return load_url_advanced(url)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN DOCUMENT DISPATCHER (unchanged from original)
# ══════════════════════════════════════════════════════════════════════════════
def load_document(file_path: str, original_name: str) -> list[Document]:
    ext = Path(original_name).suffix.lower()
    try:
        if ext in (".txt", ".html", ".htm"):
            return _load_text(file_path, original_name)
        if ext == ".eml":
            return _load_eml(file_path, original_name)
        if ext == ".pdf":
            return _load_pdf_fast(file_path, original_name)
        if ext in (".docx", ".doc"):
            return _load_docx(file_path, original_name)
        if ext in (".xlsx", ".xls"):
            return _load_excel(file_path, original_name)
        if ext in (".pptx", ".ppt"):
            return _load_pptx(file_path, original_name)
        if ext == ".csv":
            return _load_csv(file_path, original_name)
        return _docling_load(file_path, original_name)
    except Exception as exc:
        logger.error("Failed to load %s: %s", original_name, exc)
        return [Document(
            page_content=f"[Error reading {original_name}]: {exc}",
            metadata={"source": original_name, "error": str(exc)},
        )]


# ── PDF fast (pdfplumber + pypdf + Docling fallback) ──────────────────────────
def _load_pdf_fast(path: str, name: str) -> list[Document]:
    docs = []
    total_text = 0
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page_no, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        rows = []
                        for row in table:
                            cells = [str(c).strip() if c else "" for c in row]
                            rows.append(" | ".join(cells))
                        text += "\n\n" + "\n".join(rows)
                text = _normalize_units(text.strip())
                total_text += len(text)
                if text:
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": name, "page": page_no, "extraction": "pdfplumber"},
                    ))
        logger.info("pdfplumber extracted %d page(s), %d chars from %s", len(docs), total_text, name)
    except Exception as exc:
        logger.warning("pdfplumber failed for %s: %s — trying pypdf", name, exc)
        docs, total_text = _load_pdf_pypdf(path, name)

    if total_text < _MIN_TEXT_LENGTH:
        logger.info("Low text from pdfplumber (%d chars) — falling back to Docling OCR for %s", total_text, name)
        return _docling_load(path, name)

    return docs or [Document(
        page_content="(Empty PDF — no text extracted)",
        metadata={"source": name, "extraction": "pdfplumber"},
    )]


def _load_pdf_pypdf(path: str, name: str) -> tuple[list[Document], int]:
    docs = []
    total_text = 0
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        for page_no, page in enumerate(reader.pages, start=1):
            text = _normalize_units((page.extract_text() or "").strip())
            total_text += len(text)
            if text:
                docs.append(Document(
                    page_content=text,
                    metadata={"source": name, "page": page_no, "extraction": "pypdf"},
                ))
        logger.info("pypdf extracted %d page(s), %d chars from %s", len(docs), total_text, name)
    except Exception as exc:
        logger.warning("pypdf also failed for %s: %s", name, exc)
    return docs, total_text


# ── DOCX ──────────────────────────────────────────────────────────────────────
def _load_docx(path: str, name: str) -> list[Document]:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = _normalize_units("\n\n".join(paragraphs))
        if not text.strip():
            return [Document(page_content="(Empty document)", metadata={"source": name})]
        return [Document(
            page_content=text,
            metadata={"source": name, "page": 1, "extraction": "python-docx"},
        )]
    except Exception as exc:
        logger.warning("python-docx failed for %s: %s — trying Docling", name, exc)
        return _docling_load(path, name)


# ── EXCEL ─────────────────────────────────────────────────────────────────────
def _load_excel(path: str, name: str) -> list[Document]:
    try:
        import pandas as pd
        xls = pd.ExcelFile(path)
        docs = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text = f"Sheet: {sheet_name}\n\n{df.to_string(index=False)}"
            text = _normalize_units(text)
            docs.append(Document(
                page_content=text,
                metadata={"source": name, "page": sheet_name, "extraction": "pandas"},
            ))
        return docs or [Document(page_content="(Empty spreadsheet)", metadata={"source": name})]
    except Exception as exc:
        logger.warning("pandas excel failed for %s: %s — trying Docling", name, exc)
        return _docling_load(path, name)


# ── POWERPOINT ────────────────────────────────────────────────────────────────
def _load_pptx(path: str, name: str) -> list[Document]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        docs = []
        for slide_no, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    parts.append("\n".join(rows))
            text = _normalize_units("\n\n".join(parts))
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": name, "page": slide_no, "extraction": "python-pptx"},
                ))
        return docs or [Document(page_content="(Empty presentation)", metadata={"source": name})]
    except Exception as exc:
        logger.warning("python-pptx failed for %s: %s — trying Docling", name, exc)
        return _docling_load(path, name)


# ── CSV ───────────────────────────────────────────────────────────────────────
def _load_csv(path: str, name: str) -> list[Document]:
    try:
        import pandas as pd
        df = pd.read_csv(path)
        text = df.to_string(index=False)
        return [Document(
            page_content=_normalize_units(text),
            metadata={"source": name, "page": 1, "extraction": "pandas"},
        )]
    except Exception as exc:
        logger.error("CSV load failed for %s: %s", name, exc)
        return [Document(
            page_content=f"[Error reading CSV {name}]: {exc}",
            metadata={"source": name, "error": str(exc)},
        )]


# ── PLAIN TEXT / EML / HTML ─────────────────────────────────────────────────────────
def _load_text(path: str, name: str) -> list[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": name})]


def _load_eml(path: str, name: str) -> list[Document]:
    import email
    from email import policy as email_policy
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=email_policy.default)
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                body += part.get_content() or ""
            elif ct == "text/html" and not body:
                body += re.sub(r"<[^>]+>", " ", part.get_content() or "")
    else:
        body = msg.get_content() or ""
    text = (
        f"Subject: {msg.get('Subject', '')}\n"
        f"From: {msg.get('From', '')}\n"
        f"To: {msg.get('To', '')}\n"
        f"Date: {msg.get('Date', '')}\n\n"
        f"{body.strip()}"
    )
    return [Document(page_content=text, metadata={"source": name})]


# ── DOCLING FALLBACK (OCR for scanned PDFs) ───────────────────────────────────
_CONVERTER = None
def _get_converter():
    global _CONVERTER
    if _CONVERTER is None:
        from docling.document_converter import DocumentConverter
        _CONVERTER = DocumentConverter()
        logger.info("Docling DocumentConverter initialised (OCR fallback).")
    return _CONVERTER

def _docling_load(path: str, name: str) -> list[Document]:
    try:
        converter = _get_converter()
        result = converter.convert(path)
        doc = result.document
        docs = []
        pages = getattr(doc, "pages", None)
        if pages:
            for page_no, page in pages.items():
                page_items = [
                    item for item, _ in doc.iterate_items()
                    if hasattr(item, "prov") and item.prov
                    and any(p.page_no == page_no for p in item.prov)
                ]
                if page_items:
                    parts = []
                    for item in page_items:
                        if hasattr(item, "export_to_markdown"):
                            parts.append(item.export_to_markdown())
                        elif hasattr(item, "text"):
                            parts.append(item.text)
                    page_text = _normalize_units("\n\n".join(parts))
                    if page_text.strip():
                        docs.append(Document(
                            page_content=page_text,
                            metadata={"source": name, "page": page_no, "extraction": "docling-ocr"},
                        ))
        if not docs:
            markdown = _normalize_units(doc.export_to_markdown())
            if markdown.strip():
                docs.append(Document(
                    page_content=markdown,
                    metadata={"source": name, "page": 1, "extraction": "docling-ocr"},
                ))
        logger.info("Docling OCR extracted %d page(s) from %s", len(docs), name)
        return docs or [Document(
            page_content="(Empty document — no text extracted even with OCR)",
            metadata={"source": name, "extraction": "docling-ocr"},
        )]
    except Exception as exc:
        logger.error("Docling OCR also failed for %s: %s", name, exc)
        return [Document(
            page_content=f"[Extraction failed for {name}]: {exc}",
            metadata={"source": name, "error": str(exc)},
        )]